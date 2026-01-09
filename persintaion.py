#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Electronic Academy - GRADUATION PROJECT Presentation
A comprehensive, professional, and visually stunning presentation in Arabic.
Fixed: Arabic text (reshaping/bidi), modern fonts (Cairo), and project-specific content.
"""

import datetime
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not found. Please run 'pip install python-pptx'")
    exit(1)

try:
    import arabic_reshaper
    from bidi.algorithm import get_display
    HAS_RESHAPER = True
except ImportError:
    HAS_RESHAPER = False
    print("Warning: arabic-reshaper or python-bidi not found. Arabic text might look disconnected.")

# --- Theme Configuration ---
COLOR_PRIMARY = RGBColor(1, 212, 147)   # The App Primary Color (#01d493)
COLOR_DARK = RGBColor(10, 31, 59)      # Deep Navy
COLOR_SECONDARY = RGBColor(255, 255, 255) # White
COLOR_ACCENT = RGBColor(37, 99, 235)    # Link Blue
COLOR_BG_LIGHT = RGBColor(248, 250, 252)
COLOR_TEXT_MAIN = RGBColor(30, 41, 59)
COLOR_TEXT_DIM = RGBColor(100, 116, 139)

FONT_ARABIC = "Cairo"  # Modern Arabic Font
FONT_FALLBACK = "Segoe UI"

class GraduationPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333) # Widescreen 16:9
        self.prs.slide_height = Inches(7.5)
        
    def _fix_text(self, text):
        """Fixes Arabic text for libraries that don't support RTL/Shaping"""
        if not HAS_RESHAPER:
            return text
        reshaped_text = arabic_reshaper.reshape(text)
        return get_display(reshaped_text)

    def _apply_font(self, paragraph, size=18, bold=False, color=COLOR_TEXT_MAIN, align=PP_ALIGN.RIGHT):
        """Utility to apply font styles consistently"""
        paragraph.alignment = align
        run = paragraph.add_run()
        # We process text later or during assignment
        return run

    def _set_run_style(self, run, text, size=18, bold=False, color=COLOR_TEXT_MAIN):
        run.text = self._fix_text(text)
        run.font.name = FONT_ARABIC
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _add_modern_background(self, slide, dark=False):
        """Adds modern geometric background elements with App branding"""
        if dark:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_DARK
            bg.line.visible = False
            
            # Sublte Branding Line
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), self.prs.slide_width, Inches(0.2))
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_PRIMARY
            line.line.visible = False
        else:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_BG_LIGHT
            bg.line.visible = False
            
            # Sidebar Accent
            sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.1), 0, Inches(0.233), self.prs.slide_height)
            sidebar.fill.solid()
            sidebar.fill.fore_color.rgb = COLOR_PRIMARY
            sidebar.line.visible = False

    def add_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_modern_background(slide, dark=True)
        
        # Decorative Gradient-like shape
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(6), Inches(6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.fill.transparency = 0.8
        shape.line.visible = False

        # Title
        tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
        p = tx.text_frame.paragraphs[0]
        self._set_run_style(p.add_run(), "مشروع التخرج: الأكاديمية الإلكترونية", size=60, bold=True, color=COLOR_SECONDARY)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        tx2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(1))
        p2 = tx2.text_frame.paragraphs[0]
        self._set_run_style(p2.add_run(), "نظام إدارة تعلم (LMS) متطور يربط المعلم والطالب في بيئة رقمية ذكية", size=28, color=COLOR_PRIMARY)
        p2.alignment = PP_ALIGN.CENTER
        
        # Credits
        tx3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.3), Inches(0.5))
        p3 = tx3.text_frame.paragraphs[0]
        self._set_run_style(p3.add_run(), f"تاريخ العرض: {datetime.date.today().strftime('%Y/%m/%d')}", size=16, color=RGBColor(150, 150, 150))
        p3.alignment = PP_ALIGN.LEFT

    def add_section_header(self, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_modern_background(slide, dark=True)
        
        tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
        p = tx.text_frame.paragraphs[0]
        self._set_run_style(p.add_run(), title, size=54, bold=True, color=COLOR_PRIMARY)
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, content_list):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_modern_background(slide)
        
        # Header Box
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DARK
        header.line.visible = False
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0), Inches(11.5), Inches(1))
        tf = title_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        self._set_run_style(p.add_run(), title, size=36, bold=True, color=COLOR_PRIMARY)
        p.alignment = PP_ALIGN.RIGHT
        
        for i, text in enumerate(content_list):
            # Bullet box
            box = slide.shapes.add_textbox(Inches(1), Inches(1.8 + (i * 0.9)), Inches(11), Inches(0.8))
            p_text = box.text_frame.paragraphs[0]
            # Custom bullet point for RTL
            self._set_run_style(p_text.add_run(), f"{text} ●", size=22, color=COLOR_TEXT_MAIN)
            p_text.alignment = PP_ALIGN.RIGHT

    def add_tech_architecture_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_modern_background(slide)
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DARK
        header.line.visible = False
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0), Inches(11.5), Inches(1))
        self._set_run_style(title_box.text_frame.paragraphs[0].add_run(), "بنية النظام والتقنيات المستخدمة", size=36, bold=True, color=COLOR_PRIMARY)
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Layout shapes for architecture
        # Client Side
        c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(3.5), Inches(4))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = COLOR_SECONDARY
        c_box.line.color.rgb = COLOR_PRIMARY
        self._set_run_style(c_box.text_frame.paragraphs[0].add_run(), "جهة العميل (Frontend)", size=24, bold=True, color=COLOR_DARK)
        
        p = c_box.text_frame.add_paragraph()
        self._set_run_style(p.add_run(), "Next.js 14 / React", size=18)
        p = c_box.text_frame.add_paragraph()
        self._set_run_style(p.add_run(), "Tailwind CSS", size=18)
        p = c_box.text_frame.add_paragraph()
        self._set_run_style(p.add_run(), "Zustand (State)", size=18)

        # Arrow
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.7), Inches(3.5), Inches(1.5), Inches(1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_PRIMARY
        arrow.line.visible = False

        # Server Side
        s_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2), Inches(3.5), Inches(4))
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = COLOR_SECONDARY
        s_box.line.color.rgb = COLOR_PRIMARY
        self._set_run_style(s_box.text_frame.paragraphs[0].add_run(), "الخادم والبيانات (Backend)", size=24, bold=True, color=COLOR_DARK)
        
        p = s_box.text_frame.add_paragraph()
        self._set_run_style(p.add_run(), "Laravel 11 REST API", size=18)
        p = s_box.text_frame.add_paragraph()
        self._set_run_style(p.add_run(), "PostgreSQL Database", size=18)
        p = s_box.text_frame.add_paragraph()
        self._set_run_style(p.add_run(), "JWT Authentication", size=18)

    def save(self, filename):
        self.prs.save(filename)

def generate():
    pres = GraduationPresentation()
    
    # Slides content
    pres.add_title_slide()
    pres.add_section_header("نظرة عامة على المشروع")
    pres.add_content_slide("المقدمة والأهداف", [
        "تمكين المعلمين من إنشاء وبيع دوراتهم التدريبية بسهولة",
        "توفير تجربة تعليمية سلسة وتفاعلية للطلاب",
        "أتمتة العمليات المالية من خلال محفظة إلكترونية مدمجة",
        "دعم المحتوى باللغة العربية بشكل كامل وسهل الاستخدام"
    ])
    
    pres.add_section_header("المتطلبات والوظائف")
    pres.add_content_slide("أهم المتطلبات الوظيفية", [
        "نظام إدارة دورات متكامل (إضافة، تعديل، حذف)",
        "لوحة تحكم خاصة لكل من المعلم، الطالب، والمسؤول",
        "نظام اختبارات تفاعلية وتقييم تلقائي",
        "محفظة مالية تتبع المشتريات والاشتراكات"
    ])
    
    pres.add_tech_architecture_slide()
    
    pres.add_section_header("الخاتمة والنتائج")
    pres.add_content_slide("لماذا هذا المشروع؟", [
        "حل عصري يواكب التحول الرقمي في التعليم",
        "تقنيات حديثة تضمن السرعة والأمان (Next.js & SSL)",
        "قابلية للتوسع لإضافة ميزات الذكاء الاصطناعي مستقبلاً",
        "استقلالية تامة للمعلم في إدارة محتواه"
    ])
    
    output_name = "graduation-project.pptx"
    pres.save(output_name)
    print(f"DONE: Generated {output_name}")

if __name__ == "__main__":
    generate()
